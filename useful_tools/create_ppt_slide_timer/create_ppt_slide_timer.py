from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import time
import os
import win32com.client  # Used to control PowerPoint for setting transitions and exporting


def create_timer_presentation(min_timer: int) -> str:
    """
    Create a PowerPoint file including each slide with every one-second countdown.

    min_timer: int
        Number of minutes for the countdown timer.

    Return: str
        Path to the saved PowerPoint file.
    """
    # Create a PowerPoint presentation object
    presentation = Presentation()

    # Set slide width and height to 1920x1080 (in inches: 1920 px = 20 inches, 1080 px = 11.25 inches)
    presentation.slide_width = Inches(20)  # 1920px width
    presentation.slide_height = Inches(11.25)  # 1080px height

    # Start timer to track execution time
    start_time = time.time()

    # Total duration in seconds (minutes * 60 seconds)
    total_seconds = min_timer * 60

    # Loop to create slides, one for each second
    for i in range(total_seconds, -1, -1):
        # Create a new blank slide
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])  # Layout 5 is a blank slide

        # Calculate the minutes and seconds
        minutes = i // 60
        seconds = i % 60
        time_str = f"{minutes:02}:{seconds:02}"

        # Text box dimensions
        text_box_width = Inches(15)  # 3 inches width
        text_box_height = Inches(7.5)  # 1 inch height

        # Calculate center position
        left = (presentation.slide_width - text_box_width) / 2
        top = (presentation.slide_height - text_box_height) / 2

        # Add the timer text box centered on the slide
        text_box = slide.shapes.add_textbox(left, top, text_box_width, text_box_height)
        text_frame = text_box.text_frame
        text_frame.text = time_str

        # Customize the text style (font size, alignment, color)
        paragraph = text_frame.paragraphs[0]
        paragraph.font.size = Pt(450)
        paragraph.font.bold = True
        paragraph.alignment = PP_ALIGN.CENTER  # Center align the text
        paragraph.font.color.rgb = RGBColor(0, 0, 0)  # Black color for the font
        paragraph.font.name = 'EB Garamond'  # Set font to EB Garamond

    # Save the presentation
    file_name = f"{min_timer}_minute_timer_presentation.pptx"
    presentation.save(file_name)

    # End timer to track execution time
    end_time = time.time()
    print(f'Presentation created in {end_time - start_time:.2f} seconds.')
    print(f'Presentation saved as "{file_name}".')

    return file_name


def set_slide_transition_and_export(ppt_path: str, video_path: str) -> None:
    """
    Open a PowerPoint presentation, set slide transitions to 0.99 seconds, and export it as a video.
    Somehow, it can't export a video which isn't equal to the timer and each slide on PowerPoint.

    ppt_path: str
        Path to the PowerPoint presentation file.
    video_path: str
        Path to save the exported video.
    """
    # Initialize PowerPoint application
    powerpoint = win32com.client.Dispatch("PowerPoint.Application")
    powerpoint.Visible = 1  # Ensure PowerPoint is visible
    presentation = powerpoint.Presentations.Open(ppt_path)

    print('Preparation for exporting')
    # Set slide transition time to 0.99 seconds
    for slide in presentation.Slides:
        # Ensure the AdvanceOnTime is set to True
        slide.SlideShowTransition.AdvanceOnTime = True
        # Set the transition time (use this instead of directly accessing .AdvanceTime)
        slide.SlideShowTransition.AdvanceTime = float(0.99)

    print('Start exporting')
    # Export the presentation as a video
    presentation.CreateVideo(video_path, True, 0.99)

    '''
    # Save and close the presentation    
    presentation.Save()
    presentation.Close()
    powerpoint.Quit()
    '''

    print(f'Video exported as "{video_path}".')
    print('----- Make sure if the video has been exported properly. -----')


def save_slides_as_jpeg(ppt_path: str, output_folder: str) -> None:
    """
    Save each slide in the PowerPoint presentation as a JPEG image.

    ppt_path: str
        Path to the PowerPoint presentation file.
    output_folder: str
        Folder where the JPEG images will be saved.
    """
    # Initialize PowerPoint application
    powerpoint = win32com.client.Dispatch("PowerPoint.Application")
    powerpoint.Visible = 1  # Ensure PowerPoint is visible
    presentation = powerpoint.Presentations.Open(ppt_path)

    # Export each slide as a JPEG
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)

    # Loop through each slide and export as a JPEG
    for i, slide in enumerate(presentation.Slides, start=1):
        slide_path = os.path.join(output_folder, f"slide_{i}.jpg")
        slide.Export(slide_path, "JPG")
        print(f"Saved slide {i} as {slide_path}")

    # Close the presentation and quit PowerPoint
    # presentation.Close()
    # powerpoint.Quit()

    print(f'Slides have been saved on: {output_folder}')


if __name__ == '__main__':
    min_timer = int(input("Enter the number of minutes for the countdown timer: "))

    choice = int(input("Choose one of the following numbers: "
                       "1: Create PPT slide of timer and Exported as a video"
                       "2: Create PPT slide of timer and Exported as JPEG images"
                       "As Default is 2, 2 will be executed if the number entered is not 1 or 2."))

    # Get the absolute path of the current directory
    current_path = os.path.dirname(os.path.abspath(__file__))

    # Create the timer presentation
    if not os.path.exists(f"{min_timer}_minute_timer_presentation.pptx"):
        ppt_file = create_timer_presentation(min_timer)
    else:
        ppt_file = f"{min_timer}_minute_timer_presentation.pptx"
        print('The file already exists, so creating PPT file will be skipped.')

    # Ensure the correct file path is used (combine the current path with the PowerPoint file)
    ppt_file_path = os.path.join(current_path, ppt_file)

    if choice == 1:
        # Export as a video
        video_file = f"{min_timer}_minute_timer_video.mp4"
        video_file_path = os.path.join(current_path, video_file)
        set_slide_transition_and_export(ppt_file_path, video_file_path)

    else:
        # Export as JPEG images
        output_folder = os.path.join(current_path, f"{min_timer}_minute_timer_images")
        save_slides_as_jpeg(ppt_file_path, output_folder)
